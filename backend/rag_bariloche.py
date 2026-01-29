#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import logging
import json
import uuid
from datetime import datetime
from dotenv import load_dotenv
from io import BytesIO

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configuration
EMBEDDING_MODEL_NAME = os.getenv('EMBEDDING_MODEL_NAME', 'text-embedding-3-small')
TABLE_NAME = "procedimiento_rga"

# SAP AI Hub imports
EMBEDDINGS_AVAILABLE = False
ORCHESTRATION_AVAILABLE = False

try:
    from gen_ai_hub.orchestration.models.message import SystemMessage, UserMessage
    from gen_ai_hub.orchestration.models.template import Template
    from gen_ai_hub.orchestration.models.llm import LLM
    from gen_ai_hub.orchestration.models.config import OrchestrationConfig
    from gen_ai_hub.orchestration.service import OrchestrationService
    ORCHESTRATION_AVAILABLE = True
    logger.info("SAP AI Hub Orchestration cargado exitosamente")
except Exception as e:
    logger.warning(f"SAP AI Hub Orchestration no disponible: {e}")

try:
    from gen_ai_hub.proxy.core.proxy_clients import get_proxy_client
    from gen_ai_hub.proxy.native.openai import embeddings
    proxy_client = get_proxy_client("gen-ai-hub")
    EMBEDDINGS_AVAILABLE = True
    logger.info("SAP AI Hub embeddings proxy cargado exitosamente")
except Exception as e:
    logger.warning(f"SAP AI Hub embeddings proxy no disponible: {e}")

# HANA connection
HANA_AVAILABLE = False
try:
    from hana_ml import ConnectionContext
    HANA_AVAILABLE = True
    logger.info("HANA ML library cargada exitosamente")
except Exception as e:
    try:
        import hdbcli
        from hdbcli import dbapi
        HANA_AVAILABLE = True
        logger.info("HANA dbcli cargada exitosamente")
    except Exception as e2:
        logger.error(f"Bibliotecas HANA no disponibles: {e2}")

class DocumentProcessor:
    """Procesador de documentos para extraer texto"""
    
    @staticmethod
    def extract_text_from_pdf(file_content):
        """Extraer texto de archivo PDF"""
        try:
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            logger.info(f"PDF extraído: {len(text)} caracteres")
            return text
        except Exception as e:
            logger.error(f"Error extrayendo PDF: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_docx(file_content):
        """Extraer texto de documento Word"""
        try:
            from docx import Document
            doc = Document(BytesIO(file_content))
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            logger.info(f"DOCX extraído: {len(text)} caracteres")
            return text
        except Exception as e:
            logger.error(f"Error extrayendo DOCX: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_excel(file_content, filename):
        """Extraer texto de archivo Excel"""
        try:
            import pandas as pd
            excel_file = pd.ExcelFile(BytesIO(file_content))
            text = f"Archivo Excel: {filename}\n\n"
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(BytesIO(file_content), sheet_name=sheet_name)
                text += f"=== Hoja: {sheet_name} ===\n"
                text += df.to_string(index=False) + "\n\n"
            
            logger.info(f"Excel extraído: {len(excel_file.sheet_names)} hojas, {len(text)} caracteres")
            return text
        except Exception as e:
            logger.error(f"Error extrayendo Excel: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_csv(file_content):
        """Extraer texto de archivo CSV"""
        try:
            import pandas as pd
            df = pd.read_csv(BytesIO(file_content))
            text = "Datos CSV:\n\n"
            text += df.to_string(index=False)
            logger.info(f"CSV extraído: {len(df)} filas, {len(text)} caracteres")
            return text
        except Exception as e:
            logger.error(f"Error extrayendo CSV: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_pptx(file_content):
        """Extraer texto de archivo PowerPoint"""
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(file_content))
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n=== DIAPOSITIVA {slide_num} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text += " | ".join(row_text) + "\n"
            
            logger.info(f"PPTX extraído: {len(prs.slides)} diapositivas, {len(text)} caracteres")
            return text
        except Exception as e:
            logger.error(f"Error extrayendo PPTX: {e}")
            return ""
    
    @classmethod
    def extract_text(cls, file_content, filename):
        """Extraer texto del archivo basado en su extensión"""
        file_extension = filename.rsplit('.', 1)[1].lower()
        
        if file_extension == 'pdf':
            return cls.extract_text_from_pdf(file_content)
        elif file_extension == 'docx':
            return cls.extract_text_from_docx(file_content)
        elif file_extension in ['xlsx', 'xls']:
            return cls.extract_text_from_excel(file_content, filename)
        elif file_extension == 'csv':
            return cls.extract_text_from_csv(file_content)
        elif file_extension in ['pptx', 'ppt']:
            return cls.extract_text_from_pptx(file_content)
        else:
            return ""

class EmbeddingService:
    """Servicio de embeddings usando SAP AI Hub"""
    
    @staticmethod
    def get_embedding(text):
        """Obtener embedding para un texto"""
        if EMBEDDINGS_AVAILABLE:
            try:
                response = embeddings.create(model_name=EMBEDDING_MODEL_NAME, input=text)
                return response.data[0].embedding
            except Exception as e:
                logger.error(f"Error con embeddings de SAP: {e}")
                return EmbeddingService._create_fallback_embedding(text)
        else:
            return EmbeddingService._create_fallback_embedding(text)
    
    @staticmethod
    def get_embeddings(texts):
        """Generar embeddings para múltiples textos"""
        embeddings_list = []
        logger.info(f"Generando embeddings para {len(texts)} chunks")
        
        for i, text in enumerate(texts):
            embedding = EmbeddingService.get_embedding(text[:8000])
            embeddings_list.append(embedding)
            if i % 10 == 0:
                logger.info(f"Generado embedding {i+1}/{len(texts)}")
        
        return embeddings_list
    
    @staticmethod
    def _create_fallback_embedding(text):
        """Crear embedding de respaldo usando análisis de texto"""
        import hashlib
        import re
        from collections import Counter
        
        embedding = [0.0] * 1536
        
        # Hash del texto para consistencia
        text_hash = hashlib.sha256(text.encode()).hexdigest()
        for i in range(min(256, len(text_hash))):
            embedding[i] = int(text_hash[i], 16) / 15.0
        
        # Características de palabras
        words = re.findall(r'\b\w+\b', text.lower())
        word_freq = Counter(words)
        
        for i, (word, freq) in enumerate(word_freq.most_common(256)):
            if i + 256 < 1536:
                embedding[256 + i] = freq / max(len(words), 1)
        
        # Estadísticas del texto
        if len(text) > 0:
            embedding[512] = min(len(text) / 1000.0, 1.0)
            embedding[513] = min(len(words) / 100.0, 1.0)
            embedding[514] = text.count('.') / len(text)
            embedding[515] = text.count(',') / len(text)
        
        # Normalizar
        norm = sum(x * x for x in embedding) ** 0.5
        if norm > 0:
            embedding = [x / norm for x in embedding]
        
        return embedding

class HANAVectorDB:
    """Base de datos vectorial HANA"""
    
    def __init__(self):
        self.connection = None
        self.table_name = TABLE_NAME
        
    def connect(self):
        """Conectar a la base de datos HANA"""
        try:
            if HANA_AVAILABLE:
                try:
                    self.connection = ConnectionContext(
                        address=os.getenv('HANA_ADDRESS'),
                        port=int(os.getenv('HANA_PORT', 443)),
                        user=os.getenv('HANA_USER'),
                        password=os.getenv('HANA_PASSWORD'),
                        encrypt=os.getenv('HANA_ENCRYPT', 'True').lower() == 'true'
                    )
                    logger.info("Conectado a HANA usando hana_ml")
                    return True
                except:
                    # Fallback a hdbcli
                    from hdbcli import dbapi
                    self.connection = dbapi.connect(
                        address=os.getenv('HANA_ADDRESS'),
                        port=int(os.getenv('HANA_PORT', 443)),
                        user=os.getenv('HANA_USER'),
                        password=os.getenv('HANA_PASSWORD'),
                        encrypt=os.getenv('HANA_ENCRYPT', 'True').lower() == 'true'
                    )
                    logger.info("Conectado a HANA usando hdbcli")
                    return True
            return False
        except Exception as e:
            logger.error(f"Fallo al conectar a HANA: {e}")
            return False
    
    def create_table_once(self):
        """Crear la tabla solo una vez si no existe"""
        try:
            if hasattr(self.connection, 'cursor'):
                cursor = self.connection.cursor()
            else:
                cursor = self.connection.connection.cursor()
            
            # Verificar si la tabla existe
            try:
                cursor.execute(f"SELECT COUNT(*) FROM {self.table_name} WHERE 1=0")
                logger.info(f"Tabla {self.table_name} ya existe")
                cursor.close()
                return True
            except:
                # La tabla no existe, crearla
                logger.info(f"Creando tabla {self.table_name}")
                cursor.execute(f"""
                    CREATE TABLE {self.table_name} (
                        ID NVARCHAR(36) PRIMARY KEY,
                        TAG NVARCHAR(255),
                        TYPE NVARCHAR(100),
                        FILENAME NVARCHAR(255),
                        CHUNK_TEXT NCLOB,
                        CHUNK_INDEX INTEGER,
                        VECTOR_STR NCLOB,
                        VECTOR REAL_VECTOR(1536),
                        CREATED_AT TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                    )
                """)
                
                if hasattr(self.connection, 'commit'):
                    self.connection.commit()
                else:
                    self.connection.connection.commit()
                
                logger.info(f"Tabla {self.table_name} creada exitosamente")
                cursor.close()
                return True
            
        except Exception as e:
            logger.error(f"Error creando tabla: {e}")
            return False
    
    def store_document_chunks(self, filename, text_chunks, embeddings_list):
        """Almacenar chunks de documento con embeddings"""
        try:
            if not self.create_table_once():
                return False
            
            current_date = datetime.now().strftime("%Y%m%d")
            tag = f"BARILOCHE_RAG_{current_date}"
            doc_type = "PROCEDIMIENTO_DOCUMENTO"
            
            if hasattr(self.connection, 'cursor'):
                cursor = self.connection.cursor()
            else:
                cursor = self.connection.connection.cursor()
            
            stored_chunks = 0
            for i, (chunk, embedding) in enumerate(zip(text_chunks, embeddings_list)):
                chunk_id = str(uuid.uuid4())
                
                # Limpiar texto del chunk
                cleaned_chunk = chunk.replace('\n', ' ').replace('\r', ' ')
                
                sql = f"""INSERT INTO {self.table_name} 
                         (ID, TAG, TYPE, FILENAME, CHUNK_TEXT, CHUNK_INDEX, VECTOR_STR, VECTOR) 
                         VALUES (?, ?, ?, ?, ?, ?, ?, TO_REAL_VECTOR(?))"""
                
                params = (
                    chunk_id,
                    tag,
                    doc_type,
                    filename,
                    cleaned_chunk,
                    i,
                    str(embedding),
                    str(embedding)
                )
                
                cursor.execute(sql, params)
                stored_chunks += 1
            
            if hasattr(self.connection, 'commit'):
                self.connection.commit()
            else:
                self.connection.connection.commit()
            
            cursor.close()
            logger.info(f"Almacenados {stored_chunks} chunks para {filename} en {self.table_name}")
            return True
            
        except Exception as e:
            logger.error(f"Error almacenando chunks del documento: {e}")
            return False
    
    def search_similar_chunks(self, query, top_k=3):
        """Buscar chunks similares usando similitud vectorial"""
        try:
            query_embedding = EmbeddingService.get_embedding(query)
            
            if hasattr(self.connection, 'sql'):
                # Usando hana_ml
                sql = f"""
                SELECT TOP {top_k} FILENAME, CHUNK_TEXT, 
                       COSINE_SIMILARITY(VECTOR, TO_REAL_VECTOR('{query_embedding}')) as SIMILARITY
                FROM {self.table_name} 
                WHERE COSINE_SIMILARITY(VECTOR, TO_REAL_VECTOR('{query_embedding}')) > 0.1
                ORDER BY SIMILARITY DESC
                """
                try:
                    hdf = self.connection.sql(sql)
                    df_results = hdf.collect()
                    
                    similar_chunks = []
                    for _, row in df_results.iterrows():
                        similar_chunks.append({
                            'filename': row['FILENAME'],
                            'text': row['CHUNK_TEXT'][:1000],
                            'similarity': float(row['SIMILARITY'])
                        })
                except Exception as e:
                    logger.error(f"Error con consulta hana_ml: {e}")
                    # Fallback a cursor
                    cursor = self.connection.connection.cursor()
                    sql = f"""
                    SELECT TOP {top_k} FILENAME, CHUNK_TEXT, 
                           COSINE_SIMILARITY(VECTOR, TO_REAL_VECTOR(?)) as SIMILARITY
                    FROM {self.table_name} 
                    WHERE COSINE_SIMILARITY(VECTOR, TO_REAL_VECTOR(?)) > 0.1
                    ORDER BY SIMILARITY DESC
                    """
                    cursor.execute(sql, (str(query_embedding), str(query_embedding)))
                    results = cursor.fetchall()
                    
                    similar_chunks = []
                    for row in results:
                        similar_chunks.append({
                            'filename': row[0],
                            'text': row[1][:1000],
                            'similarity': float(row[2])
                        })
                    cursor.close()
            else:
                # Usando hdbcli
                cursor = self.connection.cursor()
                sql = f"""
                SELECT TOP {top_k} FILENAME, CHUNK_TEXT, 
                       COSINE_SIMILARITY(VECTOR, TO_REAL_VECTOR(?)) as SIMILARITY
                FROM {self.table_name} 
                WHERE COSINE_SIMILARITY(VECTOR, TO_REAL_VECTOR(?)) > 0.1
                ORDER BY SIMILARITY DESC
                """
                cursor.execute(sql, (str(query_embedding), str(query_embedding)))
                results = cursor.fetchall()
                
                similar_chunks = []
                for row in results:
                    similar_chunks.append({
                        'filename': row[0],
                        'text': row[1][:1000],
                        'similarity': float(row[2])
                    })
                
                cursor.close()
            
            logger.info(f"Encontrados {len(similar_chunks)} chunks similares")
            return similar_chunks
            
        except Exception as e:
            logger.error(f"Error buscando chunks similares: {e}")
            return []
    
    def get_database_stats(self):
        """Obtener estadísticas de la base de datos"""
        try:
            if hasattr(self.connection, 'cursor'):
                cursor = self.connection.cursor()
            else:
                cursor = self.connection.connection.cursor()
            
            cursor.execute(f"SELECT COUNT(*) FROM {self.table_name}")
            total_chunks = cursor.fetchone()[0]
            
            cursor.execute(f"SELECT COUNT(DISTINCT FILENAME) FROM {self.table_name}")
            unique_files = cursor.fetchone()[0]
            
            cursor.execute(f"""
                SELECT FILENAME, COUNT(*) as CHUNKS, MAX(CREATED_AT) as LAST_UPDATED
                FROM {self.table_name} 
                GROUP BY FILENAME 
                ORDER BY LAST_UPDATED DESC
            """)
            
            files_info = []
            for row in cursor.fetchall():
                files_info.append({
                    'filename': row[0],
                    'chunks': row[1],
                    'last_updated': str(row[2]) if row[2] else 'Desconocido'
                })
            
            cursor.close()
            
            return {
                'total_chunks': total_chunks,
                'unique_files': unique_files,
                'files_info': files_info
            }
            
        except Exception as e:
            logger.error(f"Error obteniendo estadísticas: {e}")
            return None
    
    def clear_all_data(self):
        """Eliminar todos los datos de la tabla"""
        try:
            if hasattr(self.connection, 'cursor'):
                cursor = self.connection.cursor()
            else:
                cursor = self.connection.connection.cursor()
            
            # Obtener el número de registros antes de eliminar
            cursor.execute(f"SELECT COUNT(*) FROM {self.table_name}")
            count_before = cursor.fetchone()[0]
            
            # Eliminar todos los registros
            cursor.execute(f"DELETE FROM {self.table_name}")
            
            if hasattr(self.connection, 'commit'):
                self.connection.commit()
            else:
                self.connection.connection.commit()
            
            cursor.close()
            
            logger.info(f"Eliminados {count_before} registros de la tabla {self.table_name}")
            return True
            
        except Exception as e:
            logger.error(f"Error eliminando datos: {e}")
            return False
    
    def close(self):
        """Cerrar conexión a la base de datos"""
        if self.connection:
            if hasattr(self.connection, 'close'):
                self.connection.close()
            else:
                self.connection.connection.close()

class RAGService:
    """Servicio principal RAG para Bariloche"""
    
    def __init__(self):
        self.vector_db = HANAVectorDB()
    
    @staticmethod
    def chunk_text(text, chunk_size=1000, overlap=100):
        """Dividir texto en chunks con superposición"""
        if not text.strip():
            return []
        
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            start += chunk_size - overlap
            if start >= len(text):
                break
        
        logger.info(f"Texto dividido en {len(chunks)} chunks")
        return chunks
    
    def process_document(self, filename, file_content):
        """Procesar documento: extraer texto, dividir y almacenar embeddings"""
        try:
            # Extraer texto
            text = DocumentProcessor.extract_text(file_content, filename)
            if not text.strip():
                return False, "No se pudo extraer texto del documento"
            
            # Dividir en chunks
            chunks = self.chunk_text(text)
            if not chunks:
                return False, "El documento no se pudo dividir en chunks"
            
            # Generar embeddings
            embeddings_list = EmbeddingService.get_embeddings(chunks)
            
            # Conectar a la base de datos y almacenar
            if not self.vector_db.connect():
                return False, "Fallo al conectar a la base de datos HANA"
            
            success = self.vector_db.store_document_chunks(filename, chunks, embeddings_list)
            
            if success:
                return True, f"Documento procesado exitosamente: {len(chunks)} chunks almacenados en {TABLE_NAME}"
            else:
                return False, "Fallo al almacenar documento en la base de datos"
                
        except Exception as e:
            logger.error(f"Error procesando documento: {e}")
            return False, f"Error procesando documento: {str(e)}"
        finally:
            self.vector_db.close()
    
    def _get_tareas_example(self):
        """Obtener un ejemplo de tareas realizadas del archivo tareas.xlsx"""
        try:
            import openpyxl
            from datetime import datetime
            
            excel_path = os.path.join('data', 'tareas.xlsx')
            if not os.path.exists(excel_path):
                return ""
            
            wb = openpyxl.load_workbook(excel_path, data_only=True)
            ws = wb.active
            
            # Leer la primera tarea como ejemplo
            tareas_ejemplo = []
            for row in ws.iter_rows(min_row=3, max_row=3, values_only=True):
                if row[0]:
                    fecha = row[0]
                    if isinstance(fecha, datetime):
                        fecha_str = fecha.strftime('%Y-%m-%d')
                    else:
                        fecha_str = str(fecha)
                    
                    tarea = {
                        'fecha': fecha_str,
                        'motivo': row[3] if row[3] else 'N/A',
                        'fenomeno': row[4] if row[4] else 'N/A',
                        'tipo_trabajo': row[5] if row[5] else 'N/A',
                        'vehiculo': row[6] if row[6] else 'N/A',
                        'equipo': row[7] if row[7] else 'N/A',
                        'urea_kg': row[8] if row[8] else 'N/A',
                        'glicol_l': row[9] if row[9] else 'N/A',
                        'prioridad_1': row[10] if row[10] else 'N/A',
                        'temperatura': row[19] if row[19] else 'N/A',
                        'humedad': row[21] if row[21] else 'N/A',
                        'viento': row[23] if row[23] else 'N/A'
                    }
                    tareas_ejemplo.append(tarea)
            
            wb.close()
            
            if tareas_ejemplo:
                tarea = tareas_ejemplo[0]
                return f"""
EJEMPLO DE TAREA REAL REALIZADA (Fecha: {tarea['fecha']}):
- Motivo de activación: {tarea['motivo']}
- Fenómeno meteorológico: {tarea['fenomeno']}
- Tipo de trabajo: {tarea['tipo_trabajo']}
- Vehículo utilizado: {tarea['vehiculo']}
- Equipo: {tarea['equipo']}
- Urea aplicada: {tarea['urea_kg']} kg
- Glicol aplicado: {tarea['glicol_l']} litros
- Prioridad de trabajo: {tarea['prioridad_1']}
- Condiciones climáticas:
  * Temperatura: {tarea['temperatura']}°C
  * Humedad: {tarea['humedad']}%
  * Viento: {tarea['viento']} km/h
"""
            return ""
            
        except Exception as e:
            logger.error(f"Error leyendo tareas.xlsx: {e}")
            return ""
    
    def _get_snow_maintenance_recommendations(self):
        """Obtener recomendaciones del procedimiento MVP1 SNOW"""
        return """
PROCEDIMIENTO MVP1 SNOW - RECOMENDACIONES PARA AVISOS DE MANTENIMIENTO:

El agente SNOW monitoreará los parámetros meteorológicos (cada 2hs) y al mismo tiempo la temperatura de pista que la obtendrá desde el MARWIS (toma el último valor medido, tener en cuenta que por ahora la temperatura de pista estará seteada en -0,7°C).

TABLA 1 - Generación Automática de Aviso 1:
Cuando las condiciones de temperatura alcancen por primera vez las de la Tabla 1 y no existe ningún otro activo en el Modo de Fallo Operativo Nieve, SNOW deberá generar automáticamente un Aviso 1.

Condiciones para Aviso 1 (Umbral de Alerta):
- Temperatura ambiente: 3°C < T ≤ 6°C
- Temperatura de rocío: Tamb. -2°C
- Temperatura de pista: 6°C > 0
- Humedad: ≥ 56%
- Viento: < 36 km/h

Clase de Aviso: Operaciones Aeropuerto
Nombre del Aviso: Umbral de Alerta
Grupo modo de fallo: Operativo Nieve
Modo de fallo: Umbral de Alerta
Ubicación técnica: RGA-LADAIR
Grupo Planificador: Operaciones

TABLA 2 - Generación de Aviso de Contingencia:
En caso de que las condiciones meteorológicas empeoren simultáneamente (tabla 2) y alcanzan los parámetros establecidos en el Umbral de Contingencia SNOW deberá generar un Aviso de Contingencia inmediatamente.

Condiciones para Aviso 2 (Umbral de Contingencia):
- Temperatura ambiente: 0°C ≤ T ≤ 3°C
- Temperatura de rocío: Tamb. -1°C
- Temperatura de pista: < 0°C
- Humedad: ≥ 63%
- Viento: < 33 km/h

Clase de Aviso: Operaciones Aeropuerto
Nombre del Aviso 2: Umbral de Contingencia
Grupo modo de fallo: Operativo Nieve
Modo de fallo: Alerta de Contingencia
Ubicación técnica: RGA-LADAIR
Grupo Planificador: Operaciones
(A futuro este aviso se convertirá en Incidencia)

TABLA 3 - Activación de Alerta de Contingencia:
Cuando las condiciones actuales meteorológicas empeoren de las establecidas en la Tabla 3 y SNOW no encuentra una lectura de Marwis en las 2hs anteriores, se activará el aviso "Alerta de Contingencia" al cambio de condiciones meteorológicas.

Condiciones para Alerta de Contingencia (sin lectura MARWIS):
- Temperatura ambiente: T ≤ 0°C
- Temperatura de rocío: Tamb. -1°C
- Temperatura de pista: < 0°C
- Humedad: ≥ 63%
- Viento: < 33 km/h

IMPORTANTE:
- Los avisos se generan automáticamente cuando se cumplen las condiciones
- El sistema SNOW monitorea cada 2 horas
- La temperatura de pista se obtiene del sistema MARWIS
- Los avisos se crean en el sistema de gestión de mantenimiento
"""

    def answer_question(self, question):
        """Responder pregunta usando RAG con información de tareas y recomendaciones SNOW"""
        try:
            # Conectar a la base de datos
            if not self.vector_db.connect():
                return "Fallo al conectar a la base de datos", []
            
            # Buscar chunks similares
            similar_chunks = self.vector_db.search_similar_chunks(question, top_k=3)
            
            if not similar_chunks:
                return "No pude encontrar información relevante en los documentos cargados.", []
            
            # Crear contexto base
            context = "\n\n".join([
                f"Documento: {chunk['filename']}\n{chunk['text']}"
                for chunk in similar_chunks
            ])
            
            # Detectar si la pregunta es sobre hielo, nieve o condiciones invernales
            keywords_hielo_nieve = ['hielo', 'nieve', 'congelamiento', 'helada', 'descongelante', 
                                    'anticongelante', 'urea', 'glicol', 'temperatura bajo cero',
                                    'frio extremo', 'pista congelada']
            
            is_winter_query = any(keyword.lower() in question.lower() for keyword in keywords_hielo_nieve)
            
            # Agregar información adicional si es consulta sobre hielo/nieve
            additional_context = ""
            if is_winter_query:
                # Agregar ejemplo de tareas
                tareas_info = self._get_tareas_example()
                if tareas_info:
                    additional_context += tareas_info
                
                # Agregar recomendaciones SNOW
                snow_recommendations = self._get_snow_maintenance_recommendations()
                additional_context += "\n" + snow_recommendations
            
            # Combinar contexto completo
            full_context = context
            if additional_context:
                full_context += "\n\n" + additional_context
            
            # Generar respuesta usando LLM
            if ORCHESTRATION_AVAILABLE:
                system_prompt = """Eres un asistente experto en procedimientos aeroportuarios. Tu tarea es proporcionar información DETALLADA y ESPECÍFICA sobre los procedimientos operativos basándote ÚNICAMENTE en los datos proporcionados.

REGLAS CRÍTICAS - NO ALUCINAR:
1. USA SOLAMENTE los datos meteorológicos EXACTOS que se te proporcionan
2. Si un valor es "undefined", "N/A" o no está disponible, NO inventes un valor - indica claramente que el dato no está disponible
3. USA la ubicación EXACTA mencionada en el análisis meteorológico (ciudad específica)
4. NO asumas valores que no están en los datos
5. Si la visibilidad es "undefined", NO digas "0 km" - di "dato no disponible"
6. Menciona EXPLÍCITAMENTE la ciudad/aeropuerto correcto del análisis meteorológico

SOBRE LOS PROCEDIMIENTOS:
1. SIEMPRE explica PASO A PASO los procedimientos encontrados en los documentos
2. NO menciones que hay procedimientos sin explicar cuáles son
3. Si encuentras procedimientos específicos, DEBES listarlos claramente con todos sus detalles
4. Incluye TODOS los pasos, requisitos, responsables y consideraciones encontradas
5. Cita el nombre exacto del documento fuente entre comillas
6. Si no encuentras procedimientos específicos, dilo claramente
7. Usa formato claro con viñetas o numeración para los procedimientos
8. Sé exhaustivo: equipos, personal, tiempos, condiciones, pero SOLO con datos reales"""

                if is_winter_query:
                    system_prompt += """

INFORMACIÓN ADICIONAL DISPONIBLE PARA CONSULTAS SOBRE HIELO/NIEVE:
- Se te proporciona un EJEMPLO REAL de tarea realizada del archivo tareas.xlsx
- Se te proporcionan las RECOMENDACIONES OFICIALES del Procedimiento MVP1 SNOW sobre cuándo crear avisos de mantenimiento
- DEBES incluir esta información en tu respuesta cuando sea relevante

CUANDO RESPONDER SOBRE HIELO/NIEVE:
1. Primero explica los procedimientos operativos del documento
2. Luego incluye el EJEMPLO DE TAREA REAL mostrando:
   - Qué se hizo en una situación similar
   - Equipos y recursos utilizados
   - Cantidades de urea/glicol aplicadas
   - Condiciones climáticas del momento
3. Después incluye las RECOMENDACIONES PARA AVISOS DE MANTENIMIENTO:
   - Explica cuándo se debe crear un Aviso 1 (Umbral de Alerta)
   - Explica cuándo se debe crear un Aviso 2 (Umbral de Contingencia)
   - Compara las condiciones actuales con las tablas de umbrales
   - Indica claramente si se debería crear un aviso según las condiciones actuales

FORMATO DE RESPUESTA ESPERADO:
- Primero: Ubicación EXACTA y resumen de datos meteorológicos REALES
- Segundo: Lista DETALLADA de los procedimientos encontrados
- Tercero: EJEMPLO DE TAREA REAL realizada (si aplica)
- Cuarto: RECOMENDACIONES PARA AVISOS DE MANTENIMIENTO con análisis de condiciones actuales (si aplica)
- Quinto: Consideraciones adicionales o alertas de seguridad
- Sexto: Fuente documental citada"""

                system_message = SystemMessage(content=system_prompt + "\n\nIMPORTANTE: NO inventes datos, NO asumas valores, NO cambies la ubicación. Sé PRECISO y HONESTO con la información disponible.")
                
                user_content = f"""CONTEXTO DE PROCEDIMIENTOS:
{full_context}

ANÁLISIS METEOROLÓGICO Y CONSULTA:
{question}

INSTRUCCIONES ESPECÍFICAS:
- Lee CUIDADOSAMENTE la ubicación exacta en el análisis meteorológico
- USA SOLO los valores meteorológicos proporcionados - si dice "undefined" o "N/A", repórtalo así
- NO inventes la visibilidad si no está disponible
- Extrae TODOS los procedimientos específicos mencionados en el contexto
- Lista CADA paso o acción mencionada en el documento
- Incluye requisitos, personal involucrado, equipos necesarios
- Si el documento menciona condiciones específicas, compáralas con los datos REALES proporcionados
- Cita textualmente las partes importantes del documento"""

                if is_winter_query and additional_context:
                    user_content += """

INFORMACIÓN ADICIONAL INCLUIDA EN EL CONTEXTO:
- EJEMPLO DE TAREA REAL del archivo tareas.xlsx (debes incluirlo en tu respuesta)
- RECOMENDACIONES OFICIALES MVP1 SNOW para avisos de mantenimiento (debes incluirlas y analizar)"""

                user_content += "\n\nResponde de forma DETALLADA, PRÁCTICA y PRECISA - sin inventar datos."
                
                user_message = UserMessage(content=user_content)
                
                llm = LLM(name="gpt-4o", version="latest")
                template = Template(messages=[system_message, user_message])
                config = OrchestrationConfig(template=template, llm=llm)
                orchestration_service = OrchestrationService(config=config)
                
                result = orchestration_service.run(template_values=[])
                answer = result.orchestration_result.choices[0].message.content.strip()
            else:
                answer = f"Basándome en los procedimientos de Bariloche, aquí está la información encontrada:\n\n{context[:500]}..."
            
            # Preparar fuentes
            sources = [
                {
                    'filename': chunk['filename'],
                    'similarity': round(chunk['similarity'], 3),
                    'text_preview': chunk['text'][:200] + "..." if len(chunk['text']) > 200 else chunk['text']
                }
                for chunk in similar_chunks
            ]
            
            return answer, sources
            
        except Exception as e:
            logger.error(f"Error respondiendo pregunta: {e}")
            return f"Error al responder la pregunta: {str(e)}", []
        finally:
            self.vector_db.close()
    
    def get_stats(self):
        """Obtener estadísticas de la base de datos"""
        try:
            if not self.vector_db.connect():
                return None
            
            stats = self.vector_db.get_database_stats()
            return stats
        except Exception as e:
            logger.error(f"Error obteniendo estadísticas: {e}")
            return None
        finally:
            self.vector_db.close()
    
    def clear_all_documents(self):
        """Limpiar todos los documentos de la base de datos"""
        try:
            if not self.vector_db.connect():
                return False, "Fallo al conectar a la base de datos HANA"
            
            success = self.vector_db.clear_all_data()
            
            if success:
                return True, "Todos los documentos han sido eliminados exitosamente"
            else:
                return False, "Error al eliminar los documentos"
                
        except Exception as e:
            logger.error(f"Error limpiando documentos: {e}")
            return False, f"Error limpiando documentos: {str(e)}"
        finally:
            self.vector_db.close()
